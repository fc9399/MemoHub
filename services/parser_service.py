# Multimodal parsing service
import os
import io
import base64
from typing import Dict, Any, Optional, List
from fastapi import UploadFile, HTTPException
import PyPDF2
import docx
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import requests
from config import settings
from services.embedding_service import embedding_service
from services.database_service import database_service

class ParserService:
    """
    Multimodal parsing service - Handles files in different formats such as text, images, audio, documents, etc.
    """
    
    def __init__(self):
        self.supported_types = {
            'pdf': self._parse_pdf,
            'docx': self._parse_docx,
            'pptx': self._parse_pptx,
            'txt': self._parse_text,
            'md': self._parse_text,
            'png': self._parse_image,
            'jpg': self._parse_image,
            'jpeg': self._parse_image,
            'gif': self._parse_image
        }
        self.max_tokens = 6000  # Set maximum token count with some margin
    
    def _split_text_into_chunks(self, text: str, max_tokens: int = None) -> List[str]:
        """Split long text into multiple chunks"""
        if max_tokens is None:
            max_tokens = self.max_tokens
        
        # More conservative token estimation: Chinese characters 1.2x, English words 1.5x
        def estimate_tokens(text: str) -> int:
            chinese_chars = len([c for c in text if '\u4e00' <= c <= '\u9fff'])
            english_words = len(text.split())
            # More conservative estimation
            return int(chinese_chars * 1.2 + english_words * 1.5)
        
        if estimate_tokens(text) <= max_tokens:
            return [text]
        
        # Split by paragraphs
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # If a single paragraph exceeds the limit, need further splitting
            if estimate_tokens(paragraph) > max_tokens:
                # Save current chunk first
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
                
                # Split overly long paragraph
                sub_chunks = self._split_long_paragraph(paragraph, max_tokens)
                chunks.extend(sub_chunks)
            elif estimate_tokens(current_chunk + paragraph) <= max_tokens:
                current_chunk += paragraph + "\n\n"
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = paragraph + "\n\n"
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_long_paragraph(self, paragraph: str, max_tokens: int) -> List[str]:
        """Split overly long paragraph"""
        def estimate_tokens(text: str) -> int:
            chinese_chars = len([c for c in text if '\u4e00' <= c <= '\u9fff'])
            english_words = len(text.split())
            return int(chinese_chars * 1.2 + english_words * 1.5)
        
        # Split by sentences
        sentences = paragraph.split('。')
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            if not sentence.strip():
                continue
                
            if estimate_tokens(current_chunk + sentence) <= max_tokens:
                current_chunk += sentence + "。"
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + "。"
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    async def parse_file(self, file: UploadFile, s3_data: Dict[str, Any], user_id: str) -> Dict[str, Any]:
        """
        Parse uploaded file and create memory unit
        
        Args:
            file: Uploaded file
            s3_data: Data returned from S3 upload
            user_id: User ID
            
        Returns:
            Dict: Parsing result and memory ID
        """
        try:
            file_extension = file.filename.split('.')[-1].lower()
            
            if file_extension not in self.supported_types:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type: {file_extension}"
                )
            
            # Reset file pointer
            await file.seek(0)
            file_content = await file.read()
            
            # Parse based on file type
            parser_func = self.supported_types[file_extension]
            parsed_content = await parser_func(file_content, file)
            
            # Generate embedding
            embedding = embedding_service.generate_embedding(
                text=parsed_content['text'],
                input_type="passage"
            )
            
            # Prepare metadata, merge additional metadata from s3_data
            metadata = {
                'original_filename': file.filename,
                'file_size': len(file_content),
                'file_extension': file_extension,
                's3_key': s3_data.get('s3_key'),
                's3_url': s3_data.get('file_url'),
                **parsed_content.get('metadata', {})
            }
            
            # If s3_data has additional metadata (e.g., Google Drive, Notion info), merge it
            if s3_data and isinstance(s3_data.get('metadata'), dict):
                metadata.update(s3_data['metadata'])
            
            # Create memory unit
            memory_id = database_service.create_memory(
                content=parsed_content['text'],
                memory_type=parsed_content['type'],
                embedding=embedding,
                user_id=user_id,
                metadata=metadata,
                source=s3_data.get('file_url'),
                summary=parsed_content.get('summary'),
                tags=parsed_content.get('tags', [])
            )
            # Process additional text chunks (if any)
            additional_memories = []
            if 'additional_chunks' in parsed_content:
                for i, chunk in enumerate(parsed_content['additional_chunks']):
                    try:
                        # Generate embedding for each chunk
                        chunk_embedding = embedding_service.generate_embedding(
                            text=chunk,
                            input_type="passage"
                        )
                        
                        # ✅ Create additional memory unit - add user_id parameter
                        chunk_memory_id = database_service.create_memory(
                            content=chunk,
                            memory_type=parsed_content['type'],
                            embedding=chunk_embedding,
                            user_id=user_id,  # ✅ Add user_id parameter
                            metadata={
                                'original_filename': file.filename,
                                'file_size': len(file_content),
                                'file_extension': file_extension,
                                's3_key': s3_data.get('s3_key'),
                                's3_url': s3_data.get('file_url'),
                                'chunk_index': i + 1,
                                'total_chunks': parsed_content['metadata'].get('total_chunks', 1),
                                'is_partial': True,
                                **parsed_content.get('metadata', {})
                            },
                            source=s3_data.get('file_url'),
                            summary=chunk[:200] + "..." if len(chunk) > 200 else chunk,
                            tags=parsed_content.get('tags', [])
                        )
                        
                        additional_memories.append({
                            'memory_id': chunk_memory_id,
                            'chunk_index': i + 1,
                            'content_preview': chunk[:100] + "..." if len(chunk) > 100 else chunk
                        })
                        
                    except Exception as e:
                        print(f"⚠️  Failed to process chunk {i + 1}: {e}")
                        import traceback
                        traceback.print_exc()
                        continue
            
            return {
                'success': True,
                'memory_id': memory_id,
                'parsed_content': parsed_content,
                'embedding_dimension': len(embedding),
                'additional_memories': additional_memories
            }
            
        except Exception as e:
            print(f"❌ File parsing failed: {e}")
            raise HTTPException(status_code=500, detail=f"Parsing failed: {str(e)}")
    
    async def _parse_pdf(self, content: bytes, file: UploadFile) -> Dict[str, Any]:
        """Parse PDF file"""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ""
            
            # First try to extract text directly
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                text += page_text + "\n"
            
            # Clean text
            text = text.strip()
            
            # If text is empty or very short, try OCR
            if not text or len(text) < 50:
                print("📄 PDF text is empty, trying OCR recognition...")
                try:
                    # Convert PDF to images
                    images = convert_from_bytes(content)
                    ocr_text = ""
                    
                    for i, image in enumerate(images):
                        print(f"🔍 OCR processing page {i+1}...")
                        # Use pytesseract for OCR
                        page_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                        ocr_text += page_text + "\n"
                    
                    if ocr_text.strip():
                        text = ocr_text.strip()
                        print(f"✅ OCR successful, extracted text length: {len(text)}")
                    else:
                        print("⚠️ OCR also failed to extract text")
                        
                except Exception as ocr_error:
                    print(f"❌ OCR failed: {ocr_error}")
                    # When OCR fails, return a message instead of empty text
                    text = f"[Scanned PDF - OCR recognition required] This PDF file is a scanned version and text cannot be extracted directly. Please install OCR dependencies or use an editable PDF file. Error details: {str(ocr_error)}"
            
            # Check text length, split into chunks if too long
            text_chunks = self._split_text_into_chunks(text)
            
            if len(text_chunks) > 1:
                # If text is chunked, return first chunk, other chunks as additional info
                return {
                    'text': text_chunks[0],
                    'type': 'document',
                    'metadata': {
                        'page_count': len(pdf_reader.pages),
                        'parser': 'PyPDF2+OCR' if len(text) > 50 else 'PyPDF2',
                        'total_chunks': len(text_chunks),
                        'chunk_index': 0,
                        'is_partial': True
                    },
                    'summary': text[:200] + "..." if len(text) > 200 else text,
                    'tags': ['pdf', 'document'],
                    'additional_chunks': text_chunks[1:]  # Other chunks
                }
            else:
                return {
                    'text': text,
                    'type': 'document',
                    'metadata': {
                        'page_count': len(pdf_reader.pages),
                        'parser': 'PyPDF2+OCR' if len(text) > 50 else 'PyPDF2'
                    },
                    'summary': text[:200] + "..." if len(text) > 200 else text,
                    'tags': ['pdf', 'document']
                }
            
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PDF parsing failed: {str(e)}")
    
    async def _parse_docx(self, content: bytes, file: UploadFile) -> Dict[str, Any]:
        """Parse DOCX file"""
        try:
            doc = docx.Document(io.BytesIO(content))
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Clean text
            text = text.strip()
            
            return {
                'text': text,
                'type': 'document',
                'metadata': {
                    'paragraph_count': len(doc.paragraphs),
                    'parser': 'python-docx'
                },
                'summary': text[:200] + "..." if len(text) > 200 else text,
                'tags': ['docx', 'document']
            }
            
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"DOCX parsing failed: {str(e)}")
    
    async def _parse_pptx(self, content: bytes, file: UploadFile) -> Dict[str, Any]:
        """Parse PPTX file"""
        try:
            try:
                from pptx import Presentation
            except ImportError:
                raise HTTPException(
                    status_code=400, 
                    detail="PPTX parsing requires 'python-pptx' package. Please install it: pip install python-pptx"
                )
            
            prs = Presentation(io.BytesIO(content))
            text = ""
            
            # Extract text from all slides
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from text boxes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    # Handle tables
                    elif hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text += cell.text + "\t"
                            text += "\n"
            
            # Clean text
            text = text.strip()
            
            return {
                'text': text,
                'type': 'document',
                'metadata': {
                    'slide_count': len(prs.slides),
                    'parser': 'python-pptx'
                },
                'summary': text[:200] + "..." if len(text) > 200 else text,
                'tags': ['pptx', 'presentation', 'document']
            }
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PPTX parsing failed: {str(e)}")
    
    async def _parse_text(self, content: bytes, file: UploadFile) -> Dict[str, Any]:
        """Parse text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            text = ""
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                text = content.decode('utf-8', errors='ignore')
            
            # Clean text
            text = text.strip()
            
            return {
                'text': text,
                'type': 'text',
                'metadata': {
                    'encoding': 'utf-8',
                    'parser': 'builtin'
                },
                'summary': text[:200] + "..." if len(text) > 200 else text,
                'tags': ['text', 'plain']
            }
            
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Text parsing failed: {str(e)}")
    
    async def _parse_image(self, content: bytes, file: UploadFile) -> Dict[str, Any]:
        """Parse image file"""
        try:
            # Open image
            image = Image.open(io.BytesIO(content))
            
            # Get image information
            width, height = image.size
            format_name = image.format
            mode = image.mode
            
            # Generate image description (simplified here, should use CLIP or OCR in practice)
            image_description = f"Image: {width}x{height} pixels, format: {format_name}, mode: {mode}"
            
            # Convert image to base64 for storage
            img_buffer = io.BytesIO()
            image.save(img_buffer, format=format_name)
            img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
            
            return {
                'text': image_description,
                'type': 'image',
                'metadata': {
                    'width': width,
                    'height': height,
                    'format': format_name,
                    'mode': mode,
                    'base64_data': img_base64,
                    'parser': 'PIL'
                },
                'summary': f"Image file: {file.filename} ({width}x{height})",
                'tags': ['image', format_name.lower()]
            }
            
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Image parsing failed: {str(e)}")
    
    async def parse_text_input(self, text: str, source: Optional[str], user_id: str) -> Dict[str, Any]:
        """Parse plain text input and create memory"""
        try:
            # Generate embedding
            embedding = embedding_service.generate_embedding(
                text=text,
                input_type="passage"
            )
            
            # ✅ Correct way - following Lin's version
            memory_id = database_service.create_memory(
                content=text,              # 1st parameter
                memory_type='text',        # 2nd parameter
                embedding=embedding,       # 3rd parameter
                user_id=user_id,          # 4th parameter ✅
                metadata={                 # 5th parameter (named)
                    'source': source or 'direct_input',
                    'parser': 'direct_text'
                },
                source=source,             # 6th parameter (named)
                summary=text[:200] + "..." if len(text) > 200 else text,  # 7th parameter
                tags=['text', 'direct_input']  # 8th parameter - you can use [] or ['text', 'direct_input']
            )
            
            return {
                'success': True,
                'memory_id': memory_id,
                'content': text,
                'embedding_dimension': len(embedding)
            }
            
        except Exception as e:
            print(f"❌ Text parsing failed: {e}")
            import traceback
            traceback.print_exc()
            raise HTTPException(status_code=500, detail=f"Text parsing failed: {str(e)}")
        
    def get_supported_types(self) -> List[str]:
        """Get supported file types"""
        return list(self.supported_types.keys())
    
    def health_check(self) -> Dict[str, Any]:
        """Health check"""
        try:
            # Check dependency libraries
            import PyPDF2
            import docx
            from PIL import Image
            
            return {
                'status': 'healthy',
                'supported_types': self.get_supported_types(),
                'dependencies': {
                    'PyPDF2': 'available',
                    'python-docx': 'available',
                    'PIL': 'available'
                }
            }
            
        except ImportError as e:
            return {
                'status': 'unhealthy',
                'error': f"Missing dependency: {str(e)}",
                'supported_types': []
            }

# Global instance
parser_service = ParserService()
